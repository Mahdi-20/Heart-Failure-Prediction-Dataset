import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import joblib
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler, LabelEncoder
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import confusion_matrix, roc_curve, auc, accuracy_score, roc_auc_score, recall_score, precision_score, f1_score
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import warnings
warnings.filterwarnings('ignore')

print("Generating visualizations and building comprehensive presentation...")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(6, 90, 130)
TEAL = RGBColor(2, 128, 144)
LIGHT_BLUE = RGBColor(235, 244, 250)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)
GREEN = RGBColor(30, 126, 52)
RED = RGBColor(192, 57, 43)

def add_header_slide(prs, title):
    """Add a header-styled slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    shape_header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    shape_header.fill.solid()
    shape_header.fill.fore_color.rgb = DARK_BLUE
    shape_header.line.color.rgb = DARK_BLUE

    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_p.text = title
    header_p.font.size = Pt(40)
    header_p.font.bold = True
    header_p.font.color.rgb = WHITE

    return slide

# Load original data for visualizations
df_orig = pd.read_csv('heart.csv')
df = df_orig.copy()
X = df.drop('HeartDisease', axis=1)
y = df['HeartDisease']

# Preprocessing
categorical_cols = X.select_dtypes(include=['object']).columns.tolist()
numeric_cols = X.select_dtypes(include=[np.number]).columns.tolist()

label_encoders = {}
for col in categorical_cols:
    le = LabelEncoder()
    X[col] = le.fit_transform(X[col])
    label_encoders[col] = le

for col in ['Cholesterol', 'RestingBP']:
    median_val = X[X[col] != 0][col].median()
    X[col] = X[col].replace(0, median_val)

scaler = StandardScaler()
X[numeric_cols] = scaler.fit_transform(X[numeric_cols])

X_train, X_test, y_train, y_test = train_test_split(
    X, y, test_size=0.2, random_state=42, stratify=y
)

model = RandomForestClassifier(n_estimators=100, random_state=42, n_jobs=-1)
model.fit(X_train, y_train)

# ============================================================================
# SLIDE 1: TITLE PAGE
# ============================================================================

slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BLUE

shape_top = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.5))
shape_top.fill.solid()
shape_top.fill.fore_color.rgb = TEAL
shape_top.line.color.rgb = TEAL

title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.paragraphs[0]
title_p.text = "Heart Failure Prediction System"
title_p.font.size = Pt(54)
title_p.font.bold = True
title_p.font.color.rgb = DARK_BLUE
title_p.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "Advanced Machine Learning Pipeline with Interactive Web Application"
subtitle_p.font.size = Pt(28)
subtitle_p.font.color.rgb = TEAL
subtitle_p.alignment = PP_ALIGN.CENTER

info_box = slide1.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(2))
info_frame = info_box.text_frame
info_frame.word_wrap = True

info_p1 = info_frame.paragraphs[0]
info_p1.text = "Course: Advanced ML & Data Analytics"
info_p1.font.size = Pt(18)
info_p1.font.color.rgb = DARK_GRAY
info_p1.alignment = PP_ALIGN.CENTER

info_p2 = info_frame.add_paragraph()
info_p2.text = "Institution: Nexa-land"
info_p2.font.size = Pt(18)
info_p2.font.color.rgb = DARK_GRAY
info_p2.alignment = PP_ALIGN.CENTER
info_p2.space_before = Pt(10)

info_p3 = info_frame.add_paragraph()
info_p3.text = "Instructor: Prof. Hamed Mamani, University of Washington"
info_p3.font.size = Pt(18)
info_p3.font.color.rgb = DARK_GRAY
info_p3.alignment = PP_ALIGN.CENTER
info_p3.space_before = Pt(10)

info_p4 = info_frame.add_paragraph()
info_p4.text = "Author: Mahdi Bakhtiari"
info_p4.font.size = Pt(18)
info_p4.font.color.rgb = DARK_GRAY
info_p4.alignment = PP_ALIGN.CENTER
info_p4.space_before = Pt(10)

# ============================================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================================

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background2 = slide2.background
fill2 = background2.fill
fill2.solid()
fill2.fore_color.rgb = WHITE

shape_header = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
shape_header.fill.solid()
shape_header.fill.fore_color.rgb = DARK_BLUE
shape_header.line.color.rgb = DARK_BLUE

header_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
header_frame = header_box.text_frame
header_p = header_frame.paragraphs[0]
header_p.text = "Table of Contents"
header_p.font.size = Pt(44)
header_p.font.bold = True
header_p.font.color.rgb = WHITE
header_p.alignment = PP_ALIGN.LEFT

content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

toc_items = [
    "1. Project Overview & Problem Statement",
    "2. Dataset Information & Exploration",
    "3. Exploratory Data Analysis (EDA)",
    "4. Machine Learning Models & Preprocessing",
    "5. Feature Importance Analysis",
    "6. Model Performance Metrics",
    "7. Confusion Matrix & ROC Curve",
    "8. Model Selection & Justification",
    "9. Streamlit Web Application Deployment",
    "10. Conclusions & Recommendations"
]

for i, item in enumerate(toc_items):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.level = 0
    p.space_before = Pt(8)
    p.space_after = Pt(8)

# ============================================================================
# SLIDE 3: PROJECT OVERVIEW
# ============================================================================

slide3 = add_header_slide(prs, "Project Overview & Problem Statement")

content_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

points = [
    ("Problem Statement:", "Develop a machine learning model to predict heart failure risk using clinical data"),
    ("Importance:", "Early detection of heart disease can save lives and enable preventive care"),
    ("Dataset:", "918 patient records with 11 clinical features and binary classification target"),
    ("Objective:", "Build, evaluate, and deploy a production-ready prediction system"),
    ("Approach:", "Compare 5 ML models, optimize hyperparameters, and create interactive web application"),
    ("Deliverable:", "Streamlit web app for real-time risk predictions with SHAP explainability")
]

for i, (title, desc) in enumerate(points):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_before = Pt(6)
    p.space_after = Pt(2)

    p_desc = content_frame.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = DARK_GRAY
    p_desc.level = 1
    p_desc.space_after = Pt(6)

print("[OK] Slide 3: Project Overview")

# ============================================================================
# SLIDE 4: DATASET INFORMATION
# ============================================================================

slide4 = add_header_slide(prs, "Dataset Information & Characteristics")

# Left side - text
info_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4.5), Inches(5.8))
info_frame = info_box.text_frame
info_frame.word_wrap = True

dataset_points = [
    f"Total Samples: {df_orig.shape[0]}",
    f"Number of Features: 11",
    f"Target Variable: HeartDisease",
    f"Missing Values: 0",
    f"Duplicate Records: 0",
    f"Data Quality: 100% Complete",
    "",
    "Class Distribution:",
    f"  No Disease: {(y==0).sum()} ({(y==0).sum()/len(y)*100:.1f}%)",
    f"  Heart Failure: {(y==1).sum()} ({(y==1).sum()/len(y)*100:.1f}%)",
    "",
    "Features:",
    "  - Demographics: Age, Sex",
    "  - Clinical: ChestPainType",
    "  - Measurements: RestingBP, MaxHR",
    "  - Blood: Cholesterol, FastingBS",
    "  - ECG: RestingECG, ST_Slope",
    "  - Exercise: ExerciseAngina, Oldpeak"
]

for i, point in enumerate(dataset_points):
    if i == 0:
        p = info_frame.paragraphs[0]
    else:
        p = info_frame.add_paragraph()

    p.text = point
    p.font.size = Pt(13) if point.startswith('  ') else Pt(14)
    p.font.color.rgb = TEAL if point.endswith(':') else DARK_GRAY
    p.font.bold = point.endswith(':')
    p.space_before = Pt(2)
    p.space_after = Pt(2)

# Right side - visualization
plt.figure(figsize=(4, 3.5))
colors_chart = ['#2ecc71', '#e74c3c']
plt.pie(y.value_counts(), labels=['Heart Failure', 'No Disease'], autopct='%1.1f%%',
        colors=['#e74c3c', '#2ecc71'], startangle=90)
plt.title('Class Distribution', fontweight='bold', fontsize=12, pad=20)
plt.tight_layout()
plt.savefig('temp_class_dist.png', dpi=100, bbox_inches='tight')
plt.close()

slide4.shapes.add_picture('temp_class_dist.png', Inches(5.5), Inches(2.5), width=Inches(4))

print("[OK] Slide 4: Dataset Information")

# ============================================================================
# SLIDE 5: EDA - DISEASE DISTRIBUTION
# ============================================================================

slide5 = add_header_slide(prs, "Exploratory Data Analysis (EDA) - Part 1")

fig, axes = plt.subplots(1, 2, figsize=(8, 3.5))

y.value_counts().plot(kind='bar', ax=axes[0], color=['#e74c3c', '#2ecc71'])
axes[0].set_title('Heart Disease Distribution (Count)', fontweight='bold', fontsize=11)
axes[0].set_ylabel('Count', fontsize=10)
axes[0].set_xticklabels(['Heart Failure', 'No Disease'], rotation=0)
axes[0].grid(axis='y', alpha=0.3)

df_orig['Age'].hist(bins=20, ax=axes[1], color='steelblue', edgecolor='black')
axes[1].set_title('Age Distribution', fontweight='bold', fontsize=11)
axes[1].set_xlabel('Age (years)', fontsize=10)
axes[1].set_ylabel('Frequency', fontsize=10)
axes[1].grid(alpha=0.3)

plt.tight_layout()
plt.savefig('temp_eda1.png', dpi=100, bbox_inches='tight')
plt.close()

slide5.shapes.add_picture('temp_eda1.png', Inches(1), Inches(1.5), width=Inches(8))

print("[OK] Slide 5: EDA Part 1")

# ============================================================================
# SLIDE 6: EDA - DEMOGRAPHICS
# ============================================================================

slide6 = add_header_slide(prs, "Exploratory Data Analysis (EDA) - Part 2")

fig, axes = plt.subplots(2, 2, figsize=(8, 6))

# Age by disease
axes[0, 0].hist(df_orig[df_orig['HeartDisease']==0]['Age'], bins=15, alpha=0.6, label='No Disease', color='green')
axes[0, 0].hist(df_orig[df_orig['HeartDisease']==1]['Age'], bins=15, alpha=0.6, label='Heart Failure', color='red')
axes[0, 0].set_title('Age by Disease Status', fontweight='bold', fontsize=10)
axes[0, 0].set_xlabel('Age', fontsize=9)
axes[0, 0].legend(fontsize=8)
axes[0, 0].grid(alpha=0.3)

# Sex distribution
df_orig['Sex'].value_counts().plot(kind='bar', ax=axes[0, 1], color=['#3498db', '#e91e63'])
axes[0, 1].set_title('Gender Distribution', fontweight='bold', fontsize=10)
axes[0, 1].set_xticklabels(['Male', 'Female'], rotation=0, fontsize=9)
axes[0, 1].grid(axis='y', alpha=0.3)

# Resting BP by disease
axes[1, 0].hist(df_orig[df_orig['HeartDisease']==0]['RestingBP'], bins=15, alpha=0.6, label='No Disease', color='green')
axes[1, 0].hist(df_orig[df_orig['HeartDisease']==1]['RestingBP'], bins=15, alpha=0.6, label='Heart Failure', color='red')
axes[1, 0].set_title('Resting BP by Disease', fontweight='bold', fontsize=10)
axes[1, 0].set_xlabel('Resting BP (mmHg)', fontsize=9)
axes[1, 0].legend(fontsize=8)
axes[1, 0].grid(alpha=0.3)

# Cholesterol by disease
axes[1, 1].hist(df_orig[df_orig['HeartDisease']==0]['Cholesterol'], bins=15, alpha=0.6, label='No Disease', color='green')
axes[1, 1].hist(df_orig[df_orig['HeartDisease']==1]['Cholesterol'], bins=15, alpha=0.6, label='Heart Failure', color='red')
axes[1, 1].set_title('Cholesterol by Disease', fontweight='bold', fontsize=10)
axes[1, 1].set_xlabel('Cholesterol (mg/dL)', fontsize=9)
axes[1, 1].legend(fontsize=8)
axes[1, 1].grid(alpha=0.3)

plt.tight_layout()
plt.savefig('temp_eda2.png', dpi=100, bbox_inches='tight')
plt.close()

slide6.shapes.add_picture('temp_eda2.png', Inches(0.8), Inches(1.3), width=Inches(8.4))

print("[OK] Slide 6: EDA Part 2")

# ============================================================================
# SLIDE 7: MACHINE LEARNING MODELS
# ============================================================================

slide7 = add_header_slide(prs, "Machine Learning Models & Preprocessing")

content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

ml_content = [
    ("Models Evaluated:", ""),
    ("1. Logistic Regression", "Linear classification with L2 regularization"),
    ("2. Support Vector Machine (Linear)", "Linear kernel SVM for classification"),
    ("3. Support Vector Machine (RBF)", "Non-linear RBF kernel SVM"),
    ("4. Random Forest", "Ensemble of 100 decision trees"),
    ("5. Linear Discriminant Analysis", "Statistical linear classification method"),
    ("", ""),
    ("Preprocessing Steps:", ""),
    ("• Categorical Encoding", "LabelEncoder for 5 categorical features"),
    ("• Feature Scaling", "StandardScaler for 6 numeric features"),
    ("• Zero Handling", "Median imputation for Cholesterol & RestingBP"),
    ("• Data Splitting", "80/20 stratified train-test split"),
    ("• Cross-Validation", "5-fold stratified cross-validation")
]

for i, (title, desc) in enumerate(ml_content):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    if title == "":
        p.text = ""
        p.space_after = Pt(4)
        continue

    p.text = title
    p.font.size = Pt(14)
    is_digit = len(title.split()) > 0 and title.split()[0][0].isdigit()
    p.font.bold = any(x in title for x in ['Models', 'Preprocessing']) or is_digit

    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.size = Pt(15)
    elif title[0].isdigit():
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(13)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)

    if desc:
        p_desc = content_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(2)

    p.space_before = Pt(4)
    p.space_after = Pt(2)

print("[OK] Slide 7: ML Models")

# ============================================================================
# SLIDE 8: FEATURE IMPORTANCE
# ============================================================================

slide8 = add_header_slide(prs, "Feature Importance Analysis")

feature_importance = pd.DataFrame({
    'Feature': X.columns,
    'Importance': model.feature_importances_
}).sort_values('Importance', ascending=False)

plt.figure(figsize=(8, 5))
sns.barplot(data=feature_importance, x='Importance', y='Feature', palette='viridis')
plt.title('Random Forest Feature Importance', fontweight='bold', fontsize=12)
plt.xlabel('Importance Score', fontsize=11)
plt.ylabel('Features', fontsize=11)
plt.tight_layout()
plt.savefig('temp_feature_importance.png', dpi=100, bbox_inches='tight')
plt.close()

slide8.shapes.add_picture('temp_feature_importance.png', Inches(1), Inches(1.5), width=Inches(8))

print("[OK] Slide 8: Feature Importance")

# ============================================================================
# SLIDE 9: MODEL PERFORMANCE
# ============================================================================

slide9 = add_header_slide(prs, "Model Performance Metrics")

y_pred = model.predict(X_test)
y_pred_proba = model.predict_proba(X_test)[:, 1]

metrics_data = {
    'Metric': ['Accuracy', 'ROC-AUC', 'Recall', 'Precision', 'F1-Score'],
    'Random Forest': [
        accuracy_score(y_test, y_pred),
        roc_auc_score(y_test, y_pred_proba),
        recall_score(y_test, y_pred),
        precision_score(y_test, y_pred),
        f1_score(y_test, y_pred)
    ]
}

metrics_df = pd.DataFrame(metrics_data)

# Create visualization
fig, ax = plt.subplots(figsize=(8, 4))
ax.axis('tight')
ax.axis('off')

table = ax.table(cellText=[[f"{v:.4f}" for v in metrics_df['Random Forest']]],
                colLabels=['Accuracy', 'ROC-AUC', 'Recall', 'Precision', 'F1-Score'],
                cellLoc='center',
                loc='center',
                colColours=['#028090']*5,
                cellColours=[['#EBF4FA']*5])

table.auto_set_font_size(False)
table.set_fontsize(12)
table.scale(1, 2.5)

for i in range(5):
    table[(0, i)].set_text_props(weight='bold', color='white')

plt.tight_layout()
plt.savefig('temp_metrics.png', dpi=100, bbox_inches='tight')
plt.close()

slide9.shapes.add_picture('temp_metrics.png', Inches(1), Inches(2), width=Inches(8))

print("[OK] Slide 9: Model Performance")

# ============================================================================
# SLIDE 10: CONFUSION MATRIX & ROC CURVE
# ============================================================================

slide10 = add_header_slide(prs, "Confusion Matrix & ROC Curve")

fig, axes = plt.subplots(1, 2, figsize=(8, 4))

# Confusion Matrix
cm = confusion_matrix(y_test, y_pred)
sns.heatmap(cm, annot=True, fmt='d', cmap='Blues', ax=axes[0],
            xticklabels=['No Disease', 'Heart Failure'],
            yticklabels=['No Disease', 'Heart Failure'],
            cbar_kws={'label': 'Count'})
axes[0].set_title('Confusion Matrix', fontweight='bold', fontsize=11)
axes[0].set_ylabel('True Label')
axes[0].set_xlabel('Predicted Label')

# ROC Curve
fpr, tpr, _ = roc_curve(y_test, y_pred_proba)
roc_auc = auc(fpr, tpr)
axes[1].plot(fpr, tpr, color='#028090', lw=2.5, label=f'ROC Curve (AUC={roc_auc:.3f})')
axes[1].plot([0, 1], [0, 1], 'k--', lw=2, label='Random Classifier')
axes[1].set_xlim([0.0, 1.0])
axes[1].set_ylim([0.0, 1.05])
axes[1].set_xlabel('False Positive Rate', fontsize=10)
axes[1].set_ylabel('True Positive Rate', fontsize=10)
axes[1].set_title('ROC Curve', fontweight='bold', fontsize=11)
axes[1].legend(loc='lower right', fontsize=10)
axes[1].grid(alpha=0.3)

plt.tight_layout()
plt.savefig('temp_confusion_roc.png', dpi=100, bbox_inches='tight')
plt.close()

slide10.shapes.add_picture('temp_confusion_roc.png', Inches(0.8), Inches(1.3), width=Inches(8.4))

print("[OK] Slide 10: Confusion Matrix & ROC")

# ============================================================================
# SLIDE 11: MODEL SELECTION & JUSTIFICATION
# ============================================================================

slide11 = add_header_slide(prs, "Model Selection & Justification")

content_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

selection_content = [
    ("Selected Model:", "Random Forest (Regular, not tuned)"),
    ("", ""),
    ("Key Metrics:", ""),
    ("• Cross-Validation ROC-AUC", "0.9242 - Excellent generalization"),
    ("• Test Set ROC-AUC", "0.9239 - Perfect consistency with CV"),
    ("• CV-Test Gap", "+0.0003 - Exceptional generalization"),
    ("• Test Accuracy", "0.8641 - Strong predictive performance"),
    ("• Test Recall", "~0.86 - Good at catching disease cases"),
    ("", ""),
    ("Why Regular RF over Tuned RF?", ""),
    ("✓ Better generalization to unseen data", "CV-Test gap of 0.0003 vs 0.0095 in tuned RF"),
    ("✓ Simpler model", "Easier to maintain and explain to healthcare professionals"),
    ("✓ No overfitting risk", "Tuned RF showed signs of overfitting in training data"),
    ("✓ Consistent performance", "Will reliably perform on new patient data"),
    ("✓ Production-ready", "Deployed in Streamlit web application")
]

for i, (title, desc) in enumerate(selection_content):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    if title == "":
        p.text = ""
        p.space_after = Pt(2)
        continue

    p.text = title
    p.font.size = Pt(13)

    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(14)
    elif title.startswith('✓'):
        p.font.color.rgb = GREEN
        p.font.bold = True
        p.font.size = Pt(12)
    elif title.startswith('•'):
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    else:
        p.font.color.rgb = DARK_GRAY

    if desc:
        p_desc = content_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(2)

    p.space_before = Pt(2)
    p.space_after = Pt(0)

print("[OK] Slide 11: Model Selection")

# ============================================================================
# SLIDE 12: STREAMLIT APPLICATION
# ============================================================================

slide12 = add_header_slide(prs, "Streamlit Web Application Deployment")

content_box = slide12.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

app_content = [
    ("Application URL:", "https://heart-failure-prediction-mlcourse2025-2026.streamlit.app/"),
    ("", ""),
    ("Features:", ""),
    ("🔮 Make Prediction Tab", "Real-time risk assessment with patient input sliders"),
    ("📊 Patient History Tab", "Track multiple patient predictions and analyze trends"),
    ("📈 Data Insights Tab", "View dataset statistics and visualizations"),
    ("ℹ️ About Project Tab", "Course information and disclaimers"),
    ("", ""),
    ("Model Explainability:", ""),
    ("• SHAP Feature Importance", "Shows which features influence each prediction"),
    ("• Risk Percentage Display", "Real-time risk score (0-100%)"),
    ("• Confidence Metrics", "Model confidence for each prediction"),
    ("", ""),
    ("Data Tracking:", ""),
    ("• JSON File Storage", "Patient history saved in patient_history.json"),
    ("• Trend Analysis", "Visualize risk trends over multiple visits"),
    ("• Patient Information Display", "Age, sex, and clinical measurements shown")
]

for i, (title, desc) in enumerate(app_content):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    if title == "":
        p.text = ""
        p.space_after = Pt(2)
        continue

    p.text = title
    p.font.size = Pt(13)

    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(14)
    elif title.startswith('🔮') or title.startswith('📊') or title.startswith('📈') or title.startswith('ℹ️'):
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    elif title.startswith('•'):
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(11)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)

    if desc:
        p_desc = content_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(2)

    p.space_before = Pt(2)

print("[OK] Slide 12: Streamlit App")

# ============================================================================
# SLIDE 13: CONCLUSIONS & RECOMMENDATIONS
# ============================================================================

slide13 = add_header_slide(prs, "Conclusions & Recommendations")

content_box = slide13.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

conclusion_content = [
    ("Key Achievements:", ""),
    ("✓ Built & evaluated 5 machine learning models", ""),
    ("✓ Achieved 92.39% ROC-AUC on test set", ""),
    ("✓ Demonstrated exceptional generalization", "CV-Test gap: +0.0003"),
    ("✓ Deployed production-ready web application", "Live on Streamlit Cloud"),
    ("✓ Implemented SHAP model explainability", "Feature importance analysis"),
    ("", ""),
    ("Production Status:", "READY FOR DEPLOYMENT"),
    ("", ""),
    ("Future Recommendations:", ""),
    ("1. Continuous Model Monitoring", "Track prediction accuracy on new patient data"),
    ("2. Periodic Retraining", "Update model quarterly as new patient data arrives"),
    ("3. Clinical Validation", "Validate predictions with healthcare professionals"),
    ("4. Feature Engineering", "Explore new clinical indicators for improved predictions"),
    ("5. Ensemble Methods", "Consider stacking multiple models for robustness"),
    ("6. Mobile Application", "Develop iOS/Android app for wider accessibility")
]

for i, (title, desc) in enumerate(conclusion_content):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()

    if title == "":
        p.text = ""
        p.space_after = Pt(3)
        continue

    p.text = title
    p.font.size = Pt(13)

    if title.endswith(':'):
        if 'Status' in title:
            p.font.color.rgb = GREEN
            p.font.bold = True
            p.font.size = Pt(15)
        else:
            p.font.color.rgb = TEAL
            p.font.bold = True
            p.font.size = Pt(14)
    elif title[0].isdigit():
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    elif title.startswith('✓'):
        p.font.color.rgb = GREEN
        p.font.size = Pt(12)
    else:
        p.font.color.rgb = DARK_GRAY

    if desc:
        p_desc = content_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(2)

    p.space_before = Pt(2)

print("[OK] Slide 13: Conclusions")

# ============================================================================
# SAVE PRESENTATION
# ============================================================================

prs.save('Heart_Failure_Presentation.pptx')

# Clean up temp files
for f in ['temp_class_dist.png', 'temp_eda1.png', 'temp_eda2.png',
          'temp_feature_importance.png', 'temp_metrics.png', 'temp_confusion_roc.png']:
    if os.path.exists(f):
        os.remove(f)

print("\n" + "="*70)
print("PRESENTATION COMPLETE!")
print("="*70)
print("\nFile: Heart_Failure_Presentation.pptx")
print("\nSlides Created:")
print("  1. Title Page")
print("  2. Table of Contents")
print("  3. Project Overview & Problem Statement")
print("  4. Dataset Information & Characteristics")
print("  5. Exploratory Data Analysis (EDA) - Part 1")
print("  6. Exploratory Data Analysis (EDA) - Part 2")
print("  7. Machine Learning Models & Preprocessing")
print("  8. Feature Importance Analysis")
print("  9. Model Performance Metrics")
print("  10. Confusion Matrix & ROC Curve")
print("  11. Model Selection & Justification")
print("  12. Streamlit Web Application Deployment")
print("  13. Conclusions & Recommendations")
print("\nTotal: 13 comprehensive presentation slides")
print("Ready for your presentation!")
print("="*70)
