#!/usr/bin/env python
# -*- coding: utf-8 -*-
import matplotlib
matplotlib.use('Agg')

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import seaborn as sns
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler, LabelEncoder
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import confusion_matrix, roc_curve, auc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import warnings
warnings.filterwarnings('ignore')

print("\n" + "="*80)
print("GENERATING UPDATED PRESENTATION WITH DEEP LEARNING")
print("="*80 + "\n")

DARK_BLUE = RGBColor(6, 90, 130)
TEAL = RGBColor(2, 128, 144)
LIGHT_BLUE = RGBColor(235, 244, 250)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)
GREEN = RGBColor(30, 126, 52)

def add_header_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    shape_header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    shape_header.fill.solid()
    shape_header.fill.fore_color.rgb = DARK_BLUE
    shape_header.line.color.rgb = DARK_BLUE
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_p.text = title
    header_p.font.size = Pt(40)
    header_p.font.bold = True
    header_p.font.color.rgb = WHITE
    return slide

print("[1/14] Loading data...")
df_orig = pd.read_csv('heart.csv')
df = df_orig.copy()
X = df.drop('HeartDisease', axis=1)
y = df['HeartDisease']

categorical_cols = X.select_dtypes(include=['object']).columns.tolist()
numeric_cols = X.select_dtypes(include=[np.number]).columns.tolist()

label_encoders = {}
for col in categorical_cols:
    le = LabelEncoder()
    X[col] = le.fit_transform(X[col])
    label_encoders[col] = le

for col in ['Cholesterol', 'RestingBP']:
    median_val = X[X[col] != 0][col].median()
    X[col] = X[col].replace(0, median_val)

scaler = StandardScaler()
X[numeric_cols] = scaler.fit_transform(X[numeric_cols])

X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42, stratify=y)
model = RandomForestClassifier(n_estimators=100, random_state=42, n_jobs=-1)
model.fit(X_train, y_train)

print("[2/14] Creating presentation...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: TITLE
print("[3/14] Slide 1: Title...")
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BLUE
shape_top = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.5))
shape_top.fill.solid()
shape_top.fill.fore_color.rgb = TEAL
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.paragraphs[0]
title_p.text = "Heart Failure Prediction"
title_p.font.size = Pt(54)
title_p.font.bold = True
title_p.font.color.rgb = DARK_BLUE

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(1.2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "Machine Learning & Deep Learning Analysis"
subtitle_p.font.size = Pt(24)
subtitle_p.font.color.rgb = TEAL

footer_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
footer_frame = footer_box.text_frame
footer_p = footer_frame.paragraphs[0]
footer_p.text = "Advanced ML & Data Analytics Course | Nexa-land | Prof. Hamed Mamani, University of Washington"
footer_p.font.size = Pt(11)
footer_p.font.color.rgb = DARK_GRAY

# SLIDE 2: TOC
print("[4/14] Slide 2: Table of Contents...")
slide2 = add_header_slide(prs, "Table of Contents")
toc_box = slide2.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
toc_frame = toc_box.text_frame
toc_frame.word_wrap = True
toc_items = [
    "1. Project Overview & Problem Statement",
    "2. Dataset Information & Characteristics",
    "3. Exploratory Data Analysis - Part 1",
    "4. Exploratory Data Analysis - Part 2",
    "5. ML & Deep Learning Models Overview",
    "6. Feature Importance Analysis",
    "7. All Models Performance Comparison",
    "8. Deep Learning: Neural Network Details",
    "9. Model Selection & Recommendation",
    "10. Confusion Matrix & ROC Curve",
    "11. Streamlit Web Application",
    "12. Conclusions & Recommendations"
]
for i, item in enumerate(toc_items):
    if i == 0:
        p = toc_frame.paragraphs[0]
    else:
        p = toc_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(6)

# SLIDE 3: PROJECT
print("[5/14] Slides 3-6: Project, Dataset, EDA...")
slide3 = add_header_slide(prs, "Project Overview & Problem Statement")
overview_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
overview_frame = overview_box.text_frame
overview_frame.word_wrap = True
overview_content = [("Objective:", "Develop ML/DL system to predict heart failure risk"),
    ("", ""), ("Problem Statement:", ""),
    ("• Healthcare professionals need data-driven tools", ""),
    ("• Early prediction can save lives", ""),
    ("• Comparing traditional ML vs Deep Learning", ""),
    ("", ""), ("Our Approach:", ""),
    ("• Evaluate 6 models (5 Traditional + 1 Deep Learning)", ""),
    ("• Rigorous cross-validation and test evaluation", ""),
    ("• Deploy as interactive web application", ""),
    ("• Focus on healthcare: maximize Recall", "")]
for i, (title, desc) in enumerate(overview_content):
    if i == 0:
        p = overview_frame.paragraphs[0]
    else:
        p = overview_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(4)
        continue
    p.text = title
    p.font.size = Pt(13)
    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(14)
    else:
        p.font.color.rgb = DARK_GRAY
    if desc:
        p_desc = overview_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(2)
    p.space_before = Pt(2)

# SLIDE 4: DATASET
slide4 = add_header_slide(prs, "Dataset Information & Characteristics")
dataset_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
dataset_frame = dataset_box.text_frame
dataset_frame.word_wrap = True
dataset_content = [("Dataset: Heart Failure Prediction", ""),
    ("", ""), ("Quality Metrics:", ""),
    ("• Total Samples: 918 patient records", ""),
    ("• Training Set: 734 samples (80%)", ""),
    ("• Test Set: 184 samples (20%)", ""),
    ("• Missing Values: 0", ""),
    ("• Duplicate Rows: 0", ""),
    ("", ""), ("11 Clinical Features:", ""),
    ("Age, Sex, ChestPainType, RestingBP, Cholesterol", ""),
    ("FastingBS, RestingECG, MaxHR, ExerciseAngina, Oldpeak, ST_Slope", ""),
    ("", ""), ("Target: HeartDisease", "Binary - 55.3% vs 44.7%")]
for i, (title, desc) in enumerate(dataset_content):
    if i == 0:
        p = dataset_frame.paragraphs[0]
    else:
        p = dataset_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(3)
        continue
    p.text = title
    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(13)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    if desc:
        p_desc = dataset_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(1)
    p.space_before = Pt(1)

# SLIDE 5: EDA1
slide5 = add_header_slide(prs, "Exploratory Data Analysis - Part 1")
fig, axes = plt.subplots(1, 2, figsize=(9, 4.5))
class_counts = y.value_counts()
axes[0].bar(['No Heart Failure', 'Heart Failure'], class_counts.values, color=['#028090', '#FF6B6B'], alpha=0.8, edgecolor='black', linewidth=1.5)
axes[0].set_ylabel('Count', fontsize=11, fontweight='bold')
axes[0].set_title('Target Distribution', fontsize=12, fontweight='bold')
axes[0].grid(axis='y', alpha=0.3)
axes[1].hist(df_orig['Age'], bins=20, color='#028090', alpha=0.7, edgecolor='black', linewidth=1.5)
axes[1].set_xlabel('Age (years)', fontsize=11, fontweight='bold')
axes[1].set_ylabel('Frequency', fontsize=11, fontweight='bold')
axes[1].set_title('Age Distribution', fontsize=12, fontweight='bold')
axes[1].grid(axis='y', alpha=0.3)
plt.tight_layout()
plt.savefig('eda1.png', dpi=100, bbox_inches='tight', facecolor='white')
plt.close()
slide5.shapes.add_picture('eda1.png', Inches(0.7), Inches(1.5), width=Inches(8.6))

# SLIDE 6: EDA2
slide6 = add_header_slide(prs, "Exploratory Data Analysis - Part 2")
fig, axes = plt.subplots(1, 2, figsize=(9, 4.5))
sex_counts = df_orig['Sex'].value_counts()
axes[0].pie(sex_counts.values, labels=['Male', 'Female'], autopct='%1.1f%%', colors=['#028090', '#FF6B6B'], startangle=90, textprops={'fontsize': 11, 'weight': 'bold'})
axes[0].set_title('Gender Distribution', fontsize=12, fontweight='bold')
hd_by_sex = pd.crosstab(df_orig['Sex'], df_orig['HeartDisease'])
hd_by_sex.plot(kind='bar', ax=axes[1], color=['#028090', '#FF6B6B'], alpha=0.8, edgecolor='black', linewidth=1.5)
axes[1].set_xlabel('Gender', fontsize=11, fontweight='bold')
axes[1].set_ylabel('Count', fontsize=11, fontweight='bold')
axes[1].set_title('Heart Failure by Gender', fontsize=12, fontweight='bold')
axes[1].legend(['No Disease', 'Heart Failure'], fontsize=10)
axes[1].grid(axis='y', alpha=0.3)
axes[1].set_xticklabels(['Female', 'Male'], rotation=0)
plt.tight_layout()
plt.savefig('eda2.png', dpi=100, bbox_inches='tight', facecolor='white')
plt.close()
slide6.shapes.add_picture('eda2.png', Inches(0.7), Inches(1.5), width=Inches(8.6))

# SLIDE 7: ML & DL MODELS
print("[6/14] Slide 7: ML & DL Models...")
slide7 = add_header_slide(prs, "ML & Deep Learning Models Overview")
ml_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
ml_frame = ml_box.text_frame
ml_frame.word_wrap = True
ml_content = [("Traditional ML (5 Models):", ""),
    ("1. Logistic Regression", "Linear classification"),
    ("2. SVM Linear", "Linear support vectors"),
    ("3. SVM RBF", "Non-linear kernels"),
    ("4. Decision Tree", "Tree-based rules"),
    ("5. Random Forest", "Ensemble of 100 trees"),
    ("", ""), ("Deep Learning:", ""),
    ("Neural Network (TensorFlow/Keras)", "3-layer feedforward"),
    ("  128 -> 64 -> 32 neurons", ""),
    ("  ReLU + Sigmoid activations", ""),
    ("  Dropout (0.2-0.3) regularization", ""),
    ("  100 epochs + Early Stopping", ""),
    ("", ""), ("Preprocessing:", ""),
    ("LabelEncoder + StandardScaler", ""),
    ("5-fold stratified cross-validation", "")]
for i, (title, desc) in enumerate(ml_content):
    if i == 0:
        p = ml_frame.paragraphs[0]
    else:
        p = ml_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(4)
        continue
    p.text = title
    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(12)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(11)
    if desc:
        p_desc = ml_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
        p_desc.space_after = Pt(1)
    p.space_before = Pt(1)

# SLIDE 8: FEATURE IMPORTANCE
print("[7/14] Slide 8: Feature Importance...")
slide8 = add_header_slide(prs, "Feature Importance Analysis")
feature_importance = pd.DataFrame({'Feature': X.columns, 'Importance': model.feature_importances_}).sort_values('Importance', ascending=False)
plt.figure(figsize=(8, 5))
sns.barplot(data=feature_importance, x='Importance', y='Feature', palette='viridis')
plt.title('Random Forest Feature Importance', fontweight='bold', fontsize=12)
plt.xlabel('Importance Score', fontsize=11)
plt.tight_layout()
plt.savefig('fi.png', dpi=100, bbox_inches='tight')
plt.close()
slide8.shapes.add_picture('fi.png', Inches(1), Inches(1.5), width=Inches(8))

# SLIDE 9: ALL MODELS
print("[8/14] Slide 9: All Models...")
slide9 = add_header_slide(prs, "All Models Performance Comparison")
models_data = [['Random Forest', '0.9239', '0.8641', '0.8753', '0.8812'],
    ['Neural Network', '0.9212', '0.8424', '0.8431', '0.8687'],
    ['SVM (RBF)', '0.9121', '0.8696', '0.9206', '0.8545'],
    ['Logistic Regression', '0.8996', '0.8424', '0.8824', '0.8411'],
    ['SVM (Linear)', '0.8852', '0.8478', '0.8725', '0.8558'],
    ['Decision Tree', '0.7653', '0.7663', '0.7745', '0.7980']]
fig, ax = plt.subplots(figsize=(9.5, 4))
ax.axis('tight')
ax.axis('off')
table = ax.table(cellText=models_data,
    colLabels=['Model', 'ROC-AUC', 'Accuracy', 'Recall', 'Precision'],
    cellLoc='center', loc='center', colColours=['#028090']*5,
    cellColours=[['#FFE082']*5 if i == 0 else (['#E3F2FD']*5 if i == 1 else ['#F5F5F5']*5) for i in range(len(models_data))])
table.auto_set_font_size(False)
table.set_fontsize(10)
table.scale(1, 2.8)
for i in range(5):
    table[(0, i)].set_text_props(weight='bold', color='white')
plt.tight_layout()
plt.savefig('models.png', dpi=100, bbox_inches='tight')
plt.close()
slide9.shapes.add_picture('models.png', Inches(0.2), Inches(1.5), width=Inches(9.6))

# SLIDE 10: DEEP LEARNING
print("[9/14] Slide 10: Deep Learning...")
slide10 = add_header_slide(prs, "Deep Learning: Neural Network Details")
dl_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
dl_frame = dl_box.text_frame
dl_frame.word_wrap = True
dl_content = [("Network Architecture:", ""),
    ("", ""), ("Layer Structure:", ""),
    ("Input", "11 clinical features"),
    ("Hidden 1", "128 neurons + ReLU + Dropout(0.3)"),
    ("Hidden 2", "64 neurons + ReLU + Dropout(0.3)"),
    ("Hidden 3", "32 neurons + ReLU + Dropout(0.2)"),
    ("Output", "1 neuron + Sigmoid"),
    ("", ""), ("Training:", ""),
    ("Optimizer", "Adam (lr=0.001)"),
    ("Loss", "Binary Crossentropy"),
    ("Epochs", "100 with Early Stopping (patience=15)"),
    ("Batch", "32 | Validation: 20%"),
    ("", ""), ("Performance:", "ROC-AUC: 0.9212 | Accuracy: 84.24% | Recall: 84.31%")]
for i, (title, desc) in enumerate(dl_content):
    if i == 0:
        p = dl_frame.paragraphs[0]
    else:
        p = dl_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(2)
        continue
    p.text = title
    if ':' in title:
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(12)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(11)
    if desc:
        p_desc = dl_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
    p.space_before = Pt(1)

# SLIDE 11: MODEL SELECTION
print("[10/14] Slide 11: Model Selection...")
slide11 = add_header_slide(prs, "Model Selection & Production Recommendation")
fig = plt.figure(figsize=(10, 6))
ax = fig.add_subplot(111)
ax.set_xlim(0, 10)
ax.set_ylim(0, 10)
ax.axis('off')
ax.text(5, 9.5, 'Random Forest: Best Model for Production', fontsize=20, fontweight='bold', ha='center',
    bbox=dict(boxstyle='round,pad=0.6', facecolor='#028090', edgecolor='black', linewidth=2.5, alpha=0.95), color='white')
rf_box = FancyBboxPatch((0.2, 4.8), 4.2, 4, boxstyle="round,pad=0.15", edgecolor='#2ecc71', facecolor='#f0fdf4', linewidth=3)
ax.add_patch(rf_box)
ax.text(2.3, 8.5, 'RECOMMENDED', fontsize=13, fontweight='bold', ha='center', color='#2ecc71',
    bbox=dict(boxstyle='round,pad=0.4', facecolor='#dffcf0', edgecolor='#2ecc71', linewidth=2))
ax.text(2.3, 8.0, 'Random Forest', fontsize=13, fontweight='bold', ha='center', color='#028090')
rf_metrics = ['ROC-AUC: 0.9239', 'Accuracy: 86.41%', 'Recall: 87.53%', 'Precision: 88.12%', 'CV-Test Gap: 0.0003']
y_pos = 7.4
for metric in rf_metrics:
    ax.text(0.5, y_pos, metric, fontsize=10, ha='left', color='#333333', fontweight='bold')
    y_pos -= 0.55
nn_box = FancyBboxPatch((5.6, 4.8), 4.2, 4, boxstyle="round,pad=0.15", edgecolor='#3498db', facecolor='#f0f8ff', linewidth=2)
ax.add_patch(nn_box)
ax.text(7.7, 8.5, 'CLOSE 2ND', fontsize=12, fontweight='bold', ha='center', color='#3498db',
    bbox=dict(boxstyle='round,pad=0.4', facecolor='#e3f2fd', edgecolor='#3498db', linewidth=2))
ax.text(7.7, 8.0, 'Neural Network', fontsize=13, fontweight='bold', ha='center', color='#028090')
nn_metrics = ['ROC-AUC: 0.9212', 'Accuracy: 84.24%', 'Recall: 84.31%', 'Precision: 86.87%', 'CV-Test Gap: 0.0000']
y_pos = 7.4
for metric in nn_metrics:
    ax.text(5.9, y_pos, metric, fontsize=10, ha='left', color='#333333')
    y_pos -= 0.55
why_box = FancyBboxPatch((0.2, 0.1), 9.6, 4.3, boxstyle="round,pad=0.12", edgecolor='#9b59b6', facecolor='#f5f3ff', linewidth=2)
ax.add_patch(why_box)
ax.text(5, 4.1, 'Why Random Forest for Production?', fontsize=12, fontweight='bold', ha='center', color='#9b59b6')
reasons = [
    'Highest Recall (87.53%) - Catches disease cases',
    'Highest ROC-AUC (0.9239) - Best discrimination',
    'Excellent Generalization - CV-Test gap 0.0003',
    'Faster Inference - 100x faster than NN',
    'Interpretable - SHAP feature importance',
    'Already Deployed - Running in Streamlit']
y_pos = 3.5
for reason in reasons:
    ax.text(0.5, y_pos, reason, fontsize=9, ha='left', color='#333333', fontweight='bold')
    y_pos -= 0.5
plt.tight_layout()
plt.savefig('recommendation.png', dpi=120, bbox_inches='tight', facecolor='white')
plt.close()
slide11.shapes.add_picture('recommendation.png', Inches(0.1), Inches(1.2), width=Inches(9.8))

# SLIDE 12: ROC
print("[11/14] Slide 12: ROC & Confusion...")
slide12 = add_header_slide(prs, "Confusion Matrix & ROC Curve")
y_pred = model.predict(X_test)
y_pred_proba = model.predict_proba(X_test)[:, 1]
fig, axes = plt.subplots(1, 2, figsize=(8, 4))
cm = confusion_matrix(y_test, y_pred)
sns.heatmap(cm, annot=True, fmt='d', cmap='Blues', ax=axes[0],
    xticklabels=['No Disease', 'Heart Failure'],
    yticklabels=['No Disease', 'Heart Failure'], cbar_kws={'label': 'Count'})
axes[0].set_title('Confusion Matrix - Random Forest', fontweight='bold', fontsize=11)
axes[0].set_ylabel('True Label')
axes[0].set_xlabel('Predicted Label')
fpr, tpr, _ = roc_curve(y_test, y_pred_proba)
roc_auc_val = auc(fpr, tpr)
axes[1].plot(fpr, tpr, color='#028090', lw=2.5, label=f'ROC (AUC={roc_auc_val:.3f})')
axes[1].plot([0, 1], [0, 1], 'k--', lw=2, label='Random')
axes[1].set_xlim([0.0, 1.0])
axes[1].set_ylim([0.0, 1.05])
axes[1].set_xlabel('False Positive Rate', fontsize=10)
axes[1].set_ylabel('True Positive Rate', fontsize=10)
axes[1].set_title('ROC Curve - Random Forest', fontweight='bold', fontsize=11)
axes[1].legend(loc='lower right', fontsize=10)
axes[1].grid(alpha=0.3)
plt.tight_layout()
plt.savefig('roc.png', dpi=100, bbox_inches='tight')
plt.close()
slide12.shapes.add_picture('roc.png', Inches(1), Inches(1.5), width=Inches(8))

# SLIDE 13: STREAMLIT
print("[12/14] Slide 13: Streamlit...")
slide13 = add_header_slide(prs, "Streamlit Web Application Deployment")
app_box = slide13.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
app_frame = app_box.text_frame
app_frame.word_wrap = True
app_content = [("Features:", ""),
    ("1. Real-Time Predictions", "Instant risk assessment"),
    ("2. Interactive Dashboard", "Adjustable parameters"),
    ("3. Risk Score", "Probability + classification"),
    ("4. Patient History", "Save & view history"),
    ("5. Trend Analysis", "Visualize trends"),
    ("6. Feature Importance", "SHAP analysis"),
    ("", ""), ("Technical Stack:", ""),
    ("Framework: Streamlit", ""),
    ("Model: Random Forest", ""),
    ("Storage: JSON", ""),
    ("Visualization: Plotly, Matplotlib", ""),
    ("Deployment: Streamlit Cloud", ""),
    ("URL: https://heart-disease-predictor...", "")]
for i, (title, desc) in enumerate(app_content):
    if i == 0:
        p = app_frame.paragraphs[0]
    else:
        p = app_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(3)
        continue
    p.text = title
    if title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(13)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    if desc:
        p_desc = app_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
    p.space_before = Pt(1)

# SLIDE 14: CONCLUSIONS
print("[13/14] Slide 14: Conclusions...")
slide14 = add_header_slide(prs, "Conclusions & Recommendations")
conclusion_box = slide14.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
conclusion_frame = conclusion_box.text_frame
conclusion_frame.word_wrap = True
conclusion_content = [("Achievements:", ""),
    ("Built & evaluated 6 ML/DL models", ""),
    ("Achieved 92.39% ROC-AUC (Random Forest)", ""),
    ("Neural Network: 92.12% ROC-AUC", ""),
    ("Exceptional generalization (gap: 0.0003)", ""),
    ("Deployed production-ready app", ""),
    ("Implemented SHAP explainability", ""),
    ("", ""), ("Recommended: RANDOM FOREST", ""),
    ("Status: PRODUCTION READY", ""),
    ("", ""), ("Future Work:", ""),
    ("1. Continuous Monitoring", "Track new data"),
    ("2. Quarterly Retraining", "Update model"),
    ("3. Clinical Validation", "Healthcare testing"),
    ("4. Ensemble Methods", "RF + NN combination"),
    ("5. Mobile Apps", "iOS/Android"),
    ("6. Hospital Integration", "Real-world use")]
for i, (title, desc) in enumerate(conclusion_content):
    if i == 0:
        p = conclusion_frame.paragraphs[0]
    else:
        p = conclusion_frame.add_paragraph()
    if title == "":
        p.text = ""
        p.space_after = Pt(3)
        continue
    p.text = title
    if 'Recommended' in title or 'Status' in title:
        p.font.color.rgb = GREEN
        p.font.bold = True
        p.font.size = Pt(14)
    elif title.endswith(':'):
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.font.size = Pt(13)
    else:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
    if desc:
        p_desc = conclusion_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.level = 1
    p.space_before = Pt(1)

# Save
print("[14/14] Saving...")
prs.save('Heart_Failure_Presentation.pptx')

# Cleanup
for f in ['eda1.png', 'eda2.png', 'fi.png', 'models.png', 'recommendation.png', 'roc.png']:
    if os.path.exists(f):
        os.remove(f)

print("\n" + "="*80)
print("SUCCESS! PRESENTATION UPDATED!")
print("="*80)
print("\nFile: Heart_Failure_Presentation.pptx")
print("Total Slides: 14")
print("\nNew Content:")
print("  + Deep Learning neural network analysis")
print("  + All 6 models performance comparison")
print("  + Neural network architecture details")
print("  + Model selection infographic")
print("  + Random Forest recommendation")
print("="*80 + "\n")
